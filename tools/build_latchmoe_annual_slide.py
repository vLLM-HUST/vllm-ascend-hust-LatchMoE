#!/usr/bin/env python3
"""Build a single-slide, editable annual-meeting overview for LatchMoE."""

from pathlib import Path
from xml.sax.saxutils import escape

import cairosvg
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "PPT"
OUT_PPTX = OUT_DIR / "LatchMoE_年会单页.pptx"
OUT_PNG = OUT_DIR / "LatchMoE_年会单页_预览.png"

SW, SH = 13.333, 7.5
FONT = "Microsoft YaHei"

BLUE = "0B43B8"
BLUE_DARK = "082E7F"
BLUE_PALE = "E8EFFB"
BLUE_LINE = "AFC3EA"
RED = "C90000"
RED_DARK = "A80B0B"
RED_PALE = "FBE8E8"
GREEN = "0A8F78"
GREEN_DARK = "076A5A"
GREEN_PALE = "E8F6F2"
INK = "161A22"
MUTED = "5C6472"
GRAY = "DDE2EA"
GRAY_DARK = "798291"
PANEL = "F9FAFC"
WHITE = "FFFFFF"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius=False):
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else kind
    sh = slide.shapes.add_shape(shape_kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    sh.line.color.rgb = rgb(line or fill)
    sh.line.width = Pt(1)
    return sh


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=16,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.04,
    font=FONT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_line(slide, x1, y1, x2, y2, color=GRAY_DARK, width=1.5):
    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(width)
    return ln


def add_chevron(slide, cx, cy, color=BLUE, scale=1.0):
    add_shape(
        slide,
        MSO_SHAPE.CHEVRON,
        cx - 0.12 * scale,
        cy - 0.10 * scale,
        0.24 * scale,
        0.20 * scale,
        color,
        color,
    )


def add_section_header(slide, label, x, y, w, color):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.36, color, color)
    add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x + w - 0.05, y + 0.04, 0.24, 0.28, color, color).rotation = 90
    add_text(slide, label, x + 0.15, y, w - 0.18, 0.36, 16, WHITE, True)


def add_goal_banner(slide):
    x, y, w, h = 0.28, 0.77, 12.77, 0.63
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, BLUE_PALE, BLUE_PALE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 1.68, h, BLUE_DARK, BLUE_DARK)
    add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x + 1.58, y + 0.05, 0.35, h - 0.10, BLUE_DARK, BLUE_DARK).rotation = 90
    add_text(slide, "研究目标", x + 0.12, y, 1.42, h, 22, WHITE, True, PP_ALIGN.CENTER)
    add_text(
        slide,
        "在受限显存下，让动态专家换入兼容 ACLGraph 稳定回放，并降低端到端推理时延",
        x + 1.94,
        y,
        w - 2.06,
        h,
        19,
        BLUE_DARK,
        True,
    )


def add_challenge_panel(slide):
    x, y, w, h = 0.28, 1.60, 2.90, 4.82
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, PANEL, "E3E7EF", True)
    add_section_header(slide, "核心挑战", x, y, 1.48, RED)
    add_text(slide, "动态专家  ×  静态图", x + 0.16, y + 0.47, w - 0.32, 0.42, 18, RED_DARK, True, PP_ALIGN.CENTER)

    cards = [
        ("1", "动态绑定", "路由专家每步变化\n图接口却要求地址稳定"),
        ("2", "容量溢出", "一次前向的活跃专家并集\n可能超过固定 slot 容量"),
        ("3", "跨流竞争", "传输流写 slot 与\n计算流读 slot 发生竞争"),
    ]
    cy = y + 1.01
    for num, title, body in cards:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.16, cy, w - 0.32, 0.94, WHITE, "E8C7C7", True)
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.28, cy + 0.18, 0.45, 0.45, RED, RED)
        add_text(slide, num, x + 0.28, cy + 0.18, 0.45, 0.45, 16, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.82, cy + 0.09, 1.54, 0.28, 14, RED_DARK, True)
        add_text(slide, body, x + 0.82, cy + 0.36, 1.58, 0.48, 11.5, INK, False, valign=MSO_ANCHOR.TOP)
        cy += 1.08

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.20, y + 4.30, w - 0.40, 0.35, RED_PALE, RED_PALE, True)
    add_text(slide, "本质：变化的逻辑状态，映射到不变的物理地址", x + 0.25, y + 4.30, w - 0.50, 0.35, 11, RED_DARK, True, PP_ALIGN.CENTER)


def add_method_panel(slide):
    x, y, w, h = 3.34, 1.60, 6.18, 4.82
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, WHITE, BLUE_LINE, True)
    add_section_header(slide, "LatchMoE 方法", x, y, 2.04, BLUE)
    add_text(slide, "逻辑专家与固定物理槽解耦", x + 2.17, y + 0.01, w - 2.30, 0.34, 15, BLUE_DARK, True)

    # Eager staging region.
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.18, y + 0.47, w - 0.36, 1.05, BLUE_PALE, BLUE_LINE, True)
    add_text(slide, "图外 Eager：动态决策与搬运", x + 0.36, y + 0.50, 2.20, 0.25, 11, BLUE_DARK, True)
    steps = [
        ("Router", "逻辑 ID"),
        ("发现活跃集", "active experts"),
        ("H2D 换入", "fixed slots"),
        ("更新映射", "log2phy"),
    ]
    sx = x + 0.32
    widths = [0.92, 1.28, 1.16, 1.12]
    for idx, ((top, bottom), bw) in enumerate(zip(steps, widths)):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, sx, y + 0.82, bw, 0.52, WHITE, BLUE_LINE, True)
        add_text(slide, top, sx + 0.04, y + 0.84, bw - 0.08, 0.23, 11, INK, True, PP_ALIGN.CENTER)
        add_text(slide, bottom, sx + 0.04, y + 1.07, bw - 0.08, 0.19, 8.5, MUTED, False, PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_chevron(slide, sx + bw + 0.19, y + 1.08, BLUE, 0.75)
        sx += bw + 0.42

    # Replay boundary.
    by = y + 1.70
    for i in range(25):
        add_line(slide, x + 0.22 + i * 0.23, by, x + 0.33 + i * 0.23, by, RED, 1.3)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 2.14, by - 0.16, 1.90, 0.32, WHITE, WHITE, True)
    add_text(slide, "ACLGraph replay boundary", x + 2.14, by - 0.16, 1.90, 0.32, 10, RED_DARK, True, PP_ALIGN.CENTER)

    # Captured region.
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.18, y + 1.88, w - 0.36, 1.68, GREEN_PALE, "A8D4C9", True)
    add_text(slide, "图内 Captured：只读取稳定地址", x + 0.36, y + 1.91, 2.30, 0.25, 11, GREEN_DARK, True)
    add_text(slide, "固定地址 Slot Bank", x + 0.47, y + 2.23, 2.22, 0.26, 11.5, GREEN_DARK, True)
    slot_x, slot_y, slot_w, slot_h = x + 0.47, y + 2.51, 2.84, 0.48
    for i in range(5):
        fill = GREEN if i in (0, 1, 3) else WHITE
        col = WHITE if fill == GREEN else GRAY_DARK
        add_shape(slide, MSO_SHAPE.RECTANGLE, slot_x + i * (slot_w / 5), slot_y, slot_w / 5 - 0.02, slot_h, fill, "9DBDB4")
        add_text(slide, f"slot{i}" if i < 4 else "…", slot_x + i * (slot_w / 5), slot_y, slot_w / 5 - 0.02, slot_h, 9.5, col, True, PP_ALIGN.CENTER)
    add_text(slide, "内容可变 · 地址不变", x + 0.47, y + 3.03, 2.84, 0.25, 10, GREEN_DARK, True, PP_ALIGN.CENTER)

    add_chevron(slide, x + 3.57, y + 2.76, GREEN, 0.85)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 3.82, y + 2.33, 0.88, 0.86, WHITE, "9DBDB4", True)
    add_text(slide, "稳定映射", x + 3.88, y + 2.40, 0.76, 0.22, 9.5, GREEN_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, "log2phy", x + 3.88, y + 2.66, 0.76, 0.30, 11, INK, True, PP_ALIGN.CENTER)
    add_chevron(slide, x + 4.93, y + 2.76, GREEN, 0.85)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 5.17, y + 2.33, 0.67, 0.86, GREEN, GREEN, True)
    add_text(slide, "Grouped\nMLP", x + 5.21, y + 2.43, 0.59, 0.62, 10.5, WHITE, True, PP_ALIGN.CENTER)

    # Three mechanism cards.
    mechs = [
        ("① 槽位虚拟化", "持久 slot + log2phy\n保持回放地址稳定", BLUE_PALE, BLUE_DARK),
        ("② 分波 Prefill", "活跃专家超容量时\n按 slot 预算分批执行", RED_PALE, RED_DARK),
        ("③ 计算保护", "LOADING → READY →\nCOMPUTING 后再复用", GREEN_PALE, GREEN_DARK),
    ]
    mx = x + 0.18
    mw = 1.84
    for title, body, fill, color in mechs:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, mx, y + 3.75, mw, 0.85, fill, fill, True)
        add_text(slide, title, mx + 0.08, y + 3.80, mw - 0.16, 0.28, 11, color, True, PP_ALIGN.CENTER)
        add_text(slide, body, mx + 0.08, y + 4.08, mw - 0.16, 0.42, 9.5, INK, False, PP_ALIGN.CENTER)
        mx += mw + 0.15


def add_metric_card(slide, x, y, title, unit, before, after, improvement, ratio, good="down"):
    w, h = 2.94, 1.15
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, WHITE, "D8DEE8", True)
    symbol = "↓" if good == "down" else "↑"
    add_text(slide, f"{title} {symbol}", x + 0.14, y + 0.08, 0.86, 0.28, 13, BLUE_DARK, True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 1.88, y + 0.07, 0.88, 0.31, GREEN_PALE, GREEN_PALE, True)
    add_text(slide, improvement, x + 1.89, y + 0.07, 0.86, 0.31, 11, GREEN_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, str(before), x + 0.14, y + 0.40, 0.90, 0.35, 16, GRAY_DARK, True, PP_ALIGN.RIGHT)
    add_text(slide, "→", x + 1.06, y + 0.40, 0.26, 0.35, 14, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, str(after), x + 1.33, y + 0.40, 0.93, 0.35, 19, GREEN_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, unit, x + 2.24, y + 0.42, 0.50, 0.31, 10, MUTED, False)
    add_text(slide, "单算子", x + 0.14, y + 0.80, 0.60, 0.22, 8.5, GRAY_DARK)
    add_text(slide, "LatchMoE", x + 1.35, y + 0.80, 0.77, 0.22, 8.5, GREEN_DARK, True)
    add_text(slide, ratio, x + 2.16, y + 0.79, 0.60, 0.24, 10.5, GREEN_DARK, True, PP_ALIGN.RIGHT)


def add_results_panel(slide):
    x, y, w, h = 9.68, 1.60, 3.37, 4.82
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, PANEL, "E3E7EF", True)
    add_section_header(slide, "实测效果", x, y, 1.48, RED)
    add_text(slide, "ShareGPT  ·  Qwen3-30B-A3B", x + 0.16, y + 0.44, w - 0.32, 0.34, 12, BLUE_DARK, True, PP_ALIGN.CENTER)
    add_metric_card(slide, x + 0.21, y + 0.88, "TTFT", "ms", "2283.2", "952.3", "−58.3%", "2.40×", "down")
    add_metric_card(slide, x + 0.21, y + 2.10, "TPOT", "ms", "192.5", "70.57", "−63.3%", "2.73×", "down")
    add_metric_card(slide, x + 0.21, y + 3.32, "吞吐", "tok/s", "4.71", "12.7", "+169.6%", "2.70×", "up")
    add_text(slide, "同一数据集与模型；数值按用户提供的实测结果计算", x + 0.24, y + 4.52, w - 0.48, 0.21, 8.5, MUTED, False, PP_ALIGN.CENTER)


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(WHITE)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.28, 0.19, 0.10, 0.45, RED, RED)
    add_text(slide, "LatchMoE：面向昇腾的图回放安全 MoE 专家卸载", 0.48, 0.12, 10.90, 0.59, 28, BLUE, True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 11.68, 0.22, 1.35, 0.37, BLUE_PALE, BLUE_PALE, True)
    add_text(slide, "年度进展 · 单页", 11.71, 0.22, 1.29, 0.37, 11, BLUE_DARK, True, PP_ALIGN.CENTER)
    add_goal_banner(slide)
    add_challenge_panel(slide)
    add_method_panel(slide)
    add_results_panel(slide)

    # Bottom conclusion ribbon.
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.28, 6.64, 12.77, 0.61, RED_PALE, RED_PALE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.28, 6.64, 1.68, 0.61, RED, RED)
    add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, 1.86, 6.69, 0.34, 0.51, RED, RED).rotation = 90
    add_text(slide, "核心结论", 0.39, 6.64, 1.43, 0.61, 21, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "首 Token 时延降低 58.3%   ·   逐 Token 时延降低 63.3%   ·   吞吐提升 2.70×", 2.20, 6.64, 10.55, 0.61, 18, RED_DARK, True, PP_ALIGN.CENTER)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)


def svg_text(x, y, text, size, color, weight=400, anchor="start", family="Noto Sans CJK SC"):
    return f'<text x="{x}" y="{y}" font-family="{family}" font-size="{size}" font-weight="{weight}" fill="#{color}" text-anchor="{anchor}">{escape(text)}</text>'


def svg_rect(x, y, w, h, fill, stroke="none", rx=0, sw=1):
    return f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="{rx}" fill="#{fill}" stroke="{stroke if stroke == "none" else "#" + stroke}" stroke-width="{sw}"/>'


def svg_multiline(x, y, lines, size, color, weight=400, anchor="start", line_h=None):
    line_h = line_h or int(size * 1.35)
    tspans = "".join(
        f'<tspan x="{x}" dy="{0 if i == 0 else line_h}">{escape(line)}</tspan>'
        for i, line in enumerate(lines)
    )
    return f'<text x="{x}" y="{y}" font-family="Noto Sans CJK SC" font-size="{size}" font-weight="{weight}" fill="#{color}" text-anchor="{anchor}">{tspans}</text>'


def build_preview():
    # A 1600x900 companion preview, proportionally matching the PPT layout.
    S = 120
    def X(v): return round(v * S)
    p = []
    p.append('<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" viewBox="0 0 1600 900">')
    p.append(svg_rect(0, 0, 1600, 900, WHITE))
    p.append(svg_rect(X(.28), X(.19), X(.10), X(.45), RED))
    p.append(svg_text(X(.48), X(.56), "LatchMoE：面向昇腾的图回放安全 MoE 专家卸载", 42, BLUE, 700))
    p.append(svg_rect(X(11.68), X(.22), X(1.35), X(.37), BLUE_PALE, rx=12))
    p.append(svg_text(X(12.355), X(.48), "年度进展 · 单页", 16, BLUE_DARK, 700, "middle"))

    # Goal banner.
    p.append(svg_rect(X(.28), X(.77), X(12.77), X(.63), BLUE_PALE))
    p.append(svg_rect(X(.28), X(.77), X(1.68), X(.63), BLUE_DARK))
    p.append(f'<polygon points="{X(1.96)},{X(.77)} {X(2.22)},{X(1.085)} {X(1.96)},{X(1.40)}" fill="#{BLUE_DARK}"/>')
    p.append(svg_text(X(1.12), X(1.20), "研究目标", 32, WHITE, 700, "middle"))
    p.append(svg_text(X(2.22), X(1.20), "在受限显存下，让动态专家换入兼容 ACLGraph 稳定回放，并降低端到端推理时延", 28, BLUE_DARK, 700))

    def panel(x, y, w, h, stroke="E3E7EF", fill=PANEL):
        p.append(svg_rect(X(x), X(y), X(w), X(h), fill, stroke, rx=17, sw=2))
    def section(x, y, w, label, color):
        p.append(svg_rect(X(x), X(y), X(w), X(.36), color))
        p.append(f'<polygon points="{X(x+w)},{X(y)} {X(x+w+.18)},{X(y+.18)} {X(x+w)},{X(y+.36)}" fill="#{color}"/>')
        p.append(svg_text(X(x+.15), X(y+.27), label, 24, WHITE, 700))

    # Challenge panel.
    panel(.28, 1.60, 2.90, 4.82)
    section(.28, 1.60, 1.48, "核心挑战", RED)
    p.append(svg_text(X(1.73), X(2.26), "动态专家  ×  静态图", 26, RED_DARK, 700, "middle"))
    cards = [
        ("1", "动态绑定", ["路由专家每步变化", "图接口却要求地址稳定"]),
        ("2", "容量溢出", ["一次前向的活跃专家并集", "可能超过固定 slot 容量"]),
        ("3", "跨流竞争", ["传输流写 slot 与", "计算流读 slot 发生竞争"]),
    ]
    cy = 2.61
    for num, title, lines in cards:
        p.append(svg_rect(X(.44), X(cy), X(2.58), X(.94), WHITE, "E8C7C7", rx=14, sw=2))
        p.append(f'<circle cx="{X(.78)}" cy="{X(cy+.405)}" r="{X(.225)}" fill="#{RED}"/>')
        p.append(svg_text(X(.78), X(cy+.48), num, 22, WHITE, 700, "middle"))
        p.append(svg_text(X(1.10), X(cy+.31), title, 19, RED_DARK, 700))
        p.append(svg_multiline(X(1.10), X(cy+.58), lines, 15, INK, 400, line_h=20))
        cy += 1.08
    p.append(svg_rect(X(.48), X(5.90), X(2.50), X(.35), RED_PALE, rx=12))
    p.append(svg_text(X(1.73), X(6.14), "本质：变化的逻辑状态，映射到不变的物理地址", 14, RED_DARK, 700, "middle"))

    # Method panel.
    panel(3.34, 1.60, 6.18, 4.82, BLUE_LINE, WHITE)
    section(3.34, 1.60, 2.04, "LatchMoE 方法", BLUE)
    p.append(svg_text(X(5.53), X(1.88), "逻辑专家与固定物理槽解耦", 21, BLUE_DARK, 700))
    p.append(svg_rect(X(3.52), X(2.07), X(5.82), X(1.05), BLUE_PALE, BLUE_LINE, rx=14, sw=2))
    p.append(svg_text(X(3.70), X(2.34), "图外 Eager：动态决策与搬运", 16, BLUE_DARK, 700))
    step_labels = [("Router", "逻辑 ID"), ("发现活跃集", "active experts"), ("H2D 换入", "fixed slots"), ("更新映射", "log2phy")]
    sx = 3.66
    bws = [.92, 1.28, 1.16, 1.12]
    for i, ((a, b), bw) in enumerate(zip(step_labels, bws)):
        p.append(svg_rect(X(sx), X(2.42), X(bw), X(.52), WHITE, BLUE_LINE, rx=9, sw=2))
        p.append(svg_text(X(sx+bw/2), X(2.65), a, 15, INK, 700, "middle"))
        p.append(svg_text(X(sx+bw/2), X(2.86), b, 11, MUTED, 400, "middle"))
        if i < 3:
            p.append(svg_text(X(sx+bw+.20), X(2.78), "›", 34, BLUE, 700, "middle"))
        sx += bw + .42
    p.append(f'<line x1="{X(3.56)}" y1="{X(3.30)}" x2="{X(9.28)}" y2="{X(3.30)}" stroke="#{RED}" stroke-width="2" stroke-dasharray="12 10"/>')
    p.append(svg_rect(X(5.48), X(3.14), X(1.90), X(.32), WHITE, rx=8))
    p.append(svg_text(X(6.43), X(3.37), "ACLGraph replay boundary", 14, RED_DARK, 700, "middle"))
    p.append(svg_rect(X(3.52), X(3.48), X(5.82), X(1.68), GREEN_PALE, "A8D4C9", rx=14, sw=2))
    p.append(svg_text(X(3.70), X(3.75), "图内 Captured：只读取稳定地址", 16, GREEN_DARK, 700))
    p.append(svg_text(X(3.81), X(4.09), "固定地址 Slot Bank", 16, GREEN_DARK, 700))
    colors = [GREEN, GREEN, WHITE, GREEN, WHITE]
    for i, c in enumerate(colors):
        p.append(svg_rect(X(3.81+i*.57), X(4.11+.38), X(.55), X(.48), c, "9DBDB4", sw=2))
        p.append(svg_text(X(4.085+i*.57), X(4.80), f"slot{i}" if i < 4 else "…", 12, WHITE if c == GREEN else GRAY_DARK, 700, "middle"))
    p.append(svg_text(X(5.23), X(5.10), "内容可变 · 地址不变", 14, GREEN_DARK, 700, "middle"))
    p.append(svg_text(X(7.03), X(4.79), "›", 38, GREEN, 700, "middle"))
    p.append(svg_rect(X(7.16), X(4.17), X(.88), X(.86), WHITE, "9DBDB4", rx=10, sw=2))
    p.append(svg_text(X(7.60), X(4.50), "稳定映射", 12, GREEN_DARK, 700, "middle"))
    p.append(svg_text(X(7.60), X(4.83), "log2phy", 15, INK, 700, "middle"))
    p.append(svg_text(X(8.28), X(4.79), "›", 38, GREEN, 700, "middle"))
    p.append(svg_rect(X(8.51), X(4.17), X(.67), X(.86), GREEN, rx=10))
    p.append(svg_multiline(X(8.845), X(4.55), ["Grouped", "MLP"], 14, WHITE, 700, "middle", 18))
    mechs = [
        ("① 槽位虚拟化", ["持久 slot + log2phy", "保持回放地址稳定"], BLUE_PALE, BLUE_DARK),
        ("② 分波 Prefill", ["活跃专家超容量时", "按 slot 预算分批执行"], RED_PALE, RED_DARK),
        ("③ 计算保护", ["LOADING → READY →", "COMPUTING 后再复用"], GREEN_PALE, GREEN_DARK),
    ]
    mx = 3.52
    for title, lines, fill, color in mechs:
        p.append(svg_rect(X(mx), X(5.35), X(1.84), X(.85), fill, rx=12))
        p.append(svg_text(X(mx+.92), X(5.64), title, 15, color, 700, "middle"))
        p.append(svg_multiline(X(mx+.92), X(5.91), lines, 12, INK, 400, "middle", 16))
        mx += 1.99

    # Results panel.
    panel(9.68, 1.60, 3.37, 4.82)
    section(9.68, 1.60, 1.48, "实测效果", RED)
    p.append(svg_text(X(11.365), X(2.30), "ShareGPT  ·  Qwen3-30B-A3B", 17, BLUE_DARK, 700, "middle"))
    metrics = [
        (2.48, "TTFT ↓", "ms", "2283.2", "952.3", "−58.3%", "2.40×"),
        (3.70, "TPOT ↓", "ms", "192.5", "70.57", "−63.3%", "2.73×"),
        (4.92, "吞吐 ↑", "tok/s", "4.71", "12.7", "+169.6%", "2.70×"),
    ]
    for my, title, unit, before, after, gain, ratio in metrics:
        p.append(svg_rect(X(9.89), X(my), X(2.94), X(1.15), WHITE, "D8DEE8", rx=14, sw=2))
        p.append(svg_text(X(10.03), X(my+.31), title, 18, BLUE_DARK, 700))
        p.append(svg_rect(X(11.77), X(my+.07), X(.88), X(.31), GREEN_PALE, rx=10))
        p.append(svg_text(X(12.21), X(my+.29), gain, 15, GREEN_DARK, 700, "middle"))
        p.append(svg_text(X(10.89), X(my+.73), before, 21, GRAY_DARK, 700, "end"))
        p.append(svg_text(X(11.09), X(my+.73), "→", 20, MUTED, 700, "middle"))
        p.append(svg_text(X(11.86), X(my+.73), after, 27, GREEN_DARK, 700, "middle"))
        p.append(svg_text(X(12.27), X(my+.71), unit, 13, MUTED, 400))
        p.append(svg_text(X(10.03), X(my+1.02), "单算子", 12, GRAY_DARK, 400))
        p.append(svg_text(X(11.24), X(my+1.02), "LatchMoE", 12, GREEN_DARK, 700))
        p.append(svg_text(X(12.64), X(my+1.02), ratio, 15, GREEN_DARK, 700, "end"))
    p.append(svg_text(X(11.365), X(6.30), "同一数据集与模型；数值按用户提供的实测结果计算", 11, MUTED, 400, "middle"))

    # Bottom ribbon.
    p.append(svg_rect(X(.28), X(6.64), X(12.77), X(.61), RED_PALE))
    p.append(svg_rect(X(.28), X(6.64), X(1.68), X(.61), RED))
    p.append(f'<polygon points="{X(1.96)},{X(6.64)} {X(2.22)},{X(6.945)} {X(1.96)},{X(7.25)}" fill="#{RED}"/>')
    p.append(svg_text(X(1.12), X(7.06), "核心结论", 29, WHITE, 700, "middle"))
    p.append(svg_text(X(7.55), X(7.06), "首 Token 时延降低 58.3%   ·   逐 Token 时延降低 63.3%   ·   吞吐提升 2.70×", 25, RED_DARK, 700, "middle"))
    p.append('</svg>')
    svg = "".join(p)
    cairosvg.svg2png(bytestring=svg.encode("utf-8"), write_to=str(OUT_PNG), output_width=1600, output_height=900)


def validate():
    prs = Presentation(OUT_PPTX)
    assert len(prs.slides) == 1
    assert prs.slide_width == Inches(SW)
    assert prs.slide_height == Inches(SH)
    text = "\n".join(sh.text for sh in prs.slides[0].shapes if hasattr(sh, "text_frame"))
    required = ["LatchMoE", "Qwen3-30B-A3B", "2283.2", "952.3", "192.5", "70.57", "4.71", "12.7"]
    missing = [item for item in required if item not in text]
    assert not missing, f"Missing slide content: {missing}"


if __name__ == "__main__":
    build_pptx()
    build_preview()
    validate()
    print(OUT_PPTX)
    print(OUT_PNG)
