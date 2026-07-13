#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate Background slide (page 4): 从动态专家路由到稳定图重放.
Left: skeleton text (1/2/3). Right: 共同机制->分岔->代价不对称 diagram.
Numbers are real (sharegpt mixed_chat, capture on vs off, 14GB, bs=1)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FONT = "Microsoft YaHei"
INK = RGBColor(0x22, 0x22, 0x22)
RED = RGBColor(0xC0, 0x1C, 0x1C)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
GREY = RGBColor(0x60, 0x60, 0x60)
BOXBLUE = RGBColor(0xE8, 0xF0, 0xFA)
BOXGREY = RGBColor(0xEF, 0xEF, 0xEF)
BOXRED = RGBColor(0xFB, 0xE8, 0xE8)
LINE = RGBColor(0x9A, 0x9A, 0x9A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def setfont(run, sz, color=INK, bold=False):
    run.font.name = FONT
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run.font._rPr
    ea = rPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        {"typeface": FONT})
    rPr.append(ea)


def tb(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tf


def para(tf, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.line_spacing = 1.05
    return p


def run(p, text, sz, color=INK, bold=False):
    r = p.add_run(); r.text = text; setfont(r, sz, color, bold); return r


# ---- title bar ----
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.82))
bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
bar.line.fill.background()
bar.shadow.inherit = False
ttf = bar.text_frame; ttf.margin_left = Inches(0.35); ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
tp = ttf.paragraphs[0]
run(tp, "Background：从动态专家路由到稳定图重放", 26, INK, True)

# ---- left skeleton text ----
LX, LW = 0.42, 6.5
tf = tb(LX, 1.02, LW, 5.9)
p = para(tf, first=True)
run(p, "① 路由驱动的动态卸载", 17, BLUE, True)
p = para(tf)
run(p, "路由每步决定激活专家；未驻留专家 CPU→NPU 按需加载。", 14)
p = para(tf)
run(p, "→ 专家驻留、权重内容、逻辑↔物理映射 持续变化。", 14, GREY)

p = para(tf); run(p, " ", 6)
p = para(tf)
run(p, "② 共同机制：索引驱动的分组矩阵乘（GPU 与 Ascend 同构）", 17, BLUE, True)
p = para(tf)
run(p, "两者都在计算前按专家排序 Token、预计算专家索引，再由分组矩阵乘按索引取权重。", 14)
p = para(tf)
run(p, "路由动态性被编码成【索引数据】——非某一平台独有。", 14, GREY)

p = para(tf); run(p, " ", 6)
p = para(tf)
run(p, "③ 图执行下的分岔：为何冲突在 Ascend 上非解不可", 17, BLUE, True)
p = para(tf)
run(p, "索引/元数据＝replay 前可更新的【数据】；专家权重地址＝capture 时【冻结】（全驻留时天然不变）。", 14)
p = para(tf)
run(p, "同一冲突，两种代价：", 14, INK, True)
p = para(tf)
run(p, "• GPU CUDA Graph：eager decode 仅轻度 host-bound，放弃图重放代价可接受（现有系统即如此）。", 13)
p = para(tf)
run(p, "• Ascend ACLGraph：每 Token ~千算子，host 追不上 → NPU 空转；放弃图重放损失过大。", 13)
p = para(tf)
run(p, "⚡ 卸载改变权重驻留与地址，与冻结的重放接口失配 → 在 Ascend 上必须正面解决。", 14, RED, True)

def box(x, y, w, h, fill, edge, lines, base_sz=12):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = edge; sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    for i, (txt, sz, col, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
        run(p, txt, sz, col, bold)
    return sh


def arrow(x1, y1, x2, y2, color=LINE, w=1.75):
    c = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    return c


RX = 7.15
box(RX + 1.05, 1.0, 3.1, 0.55, BOXGREY, LINE,
    [("每步动态路由（激活集变化）", 12, INK, True)])
arrow(RX + 2.6, 1.55, RX + 2.6, 2.05)
box(RX, 2.05, 5.6, 0.9, BOXBLUE, BLUE,
    [("共同机制：GPU & Ascend 同构", 12.5, BLUE, True),
     ("按专家排序 Token + 预计算索引 → 分组矩阵乘按索引取权重", 11, INK, False)])
tb(RX + 3.7, 2.98, 1.9, 0.3).paragraphs[0].add_run().text = ""
lbl = tb(RX, 3.0, 5.6, 0.32); lp = para(lbl, first=True); lp.alignment = PP_ALIGN.CENTER
run(lp, "图执行契约把状态一分为二", 10.5, GREY, False)
# split arrows
arrow(RX + 1.35, 3.35, RX + 1.35, 3.62)
arrow(RX + 4.25, 3.35, RX + 4.25, 3.62)
box(RX, 3.62, 2.7, 0.9, BOXGREY, LINE,
    [("索引 / 元数据", 12, INK, True),
     ("replay 前可更新 ✓", 11, RGBColor(0x2E, 0x7D, 0x32), False)])
box(RX + 2.9, 3.62, 2.7, 0.9, BOXGREY, LINE,
    [("专家权重地址", 12, INK, True),
     ("capture 时冻结（全驻留时不变）", 10.5, INK, False)])
arrow(RX + 4.25, 4.52, RX + 4.25, 4.9, RED, 2.25)
box(RX + 2.9, 4.9, 2.7, 0.62, BOXRED, RED,
    [("⚡ 卸载改变权重驻留 & 地址", 11.5, RED, True)])
arrow(RX + 4.25, 5.52, RX + 4.25, 5.87, RED, 2.25)
box(RX + 2.4, 5.87, 3.2, 0.62, BOXRED, RED,
    [("与冻结的重放接口失配 ✗", 12, RED, True)])
# callout: cost of eager fallback (number TBD — clean full-residency ACLGraph vs Eager pending)
box(RX, 5.72, 2.25, 1.15, RGBColor(0xFF, 0xFB, 0xE6), RGBColor(0xC9, 0x9A, 0x00),
    [("图重放价值（待补）", 9.5, GREY, True),
     ("全驻留 ACLGraph vs Eager", 9.5, INK, True),
     ("TPOT / 吞吐 差距：X×", 10, RED, True),
     ("（双卡全驻留实验）", 8.5, GREY, False)])

prs.save("/root/vllm-moe-offload-ascend/docs/PPT/page4_background.pptx")
print("saved full page4")
