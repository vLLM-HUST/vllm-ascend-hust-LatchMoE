from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("SEW-Offload_年会一页PPT.pptx")


BLUE = RGBColor(0x00, 0x45, 0xAD)
NAVY = RGBColor(0x0A, 0x35, 0x73)
RED = RGBColor(0xC9, 0x00, 0x00)
PALE_BLUE = RGBColor(0xE9, 0xEF, 0xF8)
PALE_RED = RGBColor(0xFA, 0xE3, 0xE3)
GREEN = RGBColor(0x2D, 0x8B, 0x6B)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xEE)
ORANGE = RGBColor(0xF0, 0x9A, 0x1A)
LIGHT_ORANGE = RGBColor(0xFE, 0xF4, 0xE4)
GREY = RGBColor(0xF3, 0xF5, 0xF7)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MID = RGBColor(0x5C, 0x6B, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, font="Microsoft YaHei", valign=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line if line else fill
    shp.line.width = Pt(1.0)
    if radius:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp


def add_triangle(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    return shp


def add_band(slide, y, label, body, label_fill, body_fill, h=0.64,
             body_size=20, body_bold=True):
    add_rect(slide, 0.30, y, 2.05, h, label_fill)
    add_triangle(slide, 2.12, y, 0.32, h, label_fill)
    add_rect(slide, 2.42, y, 10.48, h, body_fill)
    add_textbox(slide, 0.38, y + 0.02, 1.75, h - 0.04, label, 22, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, 2.55, y + 0.02, 10.18, h - 0.04, body, body_size, body_bold, DARK)


def add_arrow(slide, x1, y1, x2, y2, color=MID, width=2.0, begin=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    if begin:
        conn.line.begin_arrowhead = True
    return conn


def add_pill(slide, x, y, w, h, text, fill, line=None, size=12, bold=True, color=DARK):
    add_rect(slide, x, y, w, h, fill, line or fill, radius=True)
    add_textbox(slide, x + 0.03, y, w - 0.06, h, text, size, bold, color, PP_ALIGN.CENTER)


def add_card(slide, x, y, w, h, title, value, subtitle, fill, accent):
    add_rect(slide, x, y, w, h, fill, RGBColor(0xD7, 0xDF, 0xEA), radius=True)
    add_rect(slide, x, y, 0.08, h, accent, accent, radius=False)
    add_textbox(slide, x + 0.15, y + 0.08, w - 0.25, 0.22, title, 9.5, True, MID)
    add_textbox(slide, x + 0.15, y + 0.32, w - 0.25, 0.36, value, 17, True, accent)
    add_textbox(slide, x + 0.15, y + 0.70, w - 0.25, 0.28, subtitle, 8.5, False, MID)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.333, 7.5, WHITE)

    add_textbox(
        slide,
        0.24,
        0.14,
        12.85,
        0.52,
        "任务3-4：MoE 专家卸载与昇腾图执行适配",
        29,
        True,
        BLUE,
    )

    add_band(
        slide,
        0.86,
        "研究目标",
        "在单卡 Ascend HBM 受限场景下，让动态 MoE Expert 卸载仍能稳定 ACLGraph 重放",
        NAVY,
        PALE_BLUE,
        h=0.58,
        body_size=20,
    )

    add_band(
        slide,
        1.56,
        "方法1",
        "固定物理 Slot + 持久 log2phy：把“动态 Expert 路由”变成“稳定地址重放”",
        RED,
        PALE_RED,
        h=0.68,
        body_size=19,
    )
    add_textbox(
        slide,
        2.58,
        2.02,
        9.95,
        0.22,
        "Host Expert Store 按需加载；Grouped MLP 只看到固定 NPU Slot，不感知专家搬移",
        10.5,
        False,
        MID,
    )

    add_band(
        slide,
        2.34,
        "方法2",
        "B2 波次化 Prefill + COMPUTING 保护：专家数超过槽位仍不改语义、不覆盖在算权重",
        RED,
        PALE_RED,
        h=0.68,
        body_size=18,
    )

    # Left failure example.
    add_rect(slide, 0.48, 3.20, 3.02, 2.25, RGBColor(0xFF, 0xF7, 0xF7), RGBColor(0xE9, 0xB8, 0xB8), radius=True)
    add_textbox(slide, 0.62, 3.28, 2.75, 0.28, "现有路径：动态搬移撞上图边界", 12.5, True, RED, PP_ALIGN.CENTER)
    add_pill(slide, 0.72, 3.75, 0.82, 0.36, "Router", PALE_BLUE, RGBColor(0xB9, 0xC7, 0xD8), 10)
    add_pill(slide, 1.86, 3.75, 1.05, 0.36, "H2D Copy", PALE_RED, RGBColor(0xE4, 0xB0, 0xB0), 10)
    add_arrow(slide, 1.55, 3.93, 1.86, 3.93, RED, 2.2)
    add_arrow(slide, 2.90, 3.93, 3.24, 3.93, RED, 2.2)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.02), Inches(3.62), Inches(3.02), Inches(4.65))
    line.line.color.rgb = RED
    line.line.width = Pt(2)
    line.line.dash_style = 4
    add_textbox(slide, 2.68, 4.68, 0.72, 0.25, "ACLGraph\n边界", 8.5, True, RED, PP_ALIGN.CENTER)
    add_rect(slide, 0.76, 4.48, 2.28, 0.46, WHITE, RGBColor(0xDE, 0xA0, 0xA0), radius=True)
    add_textbox(slide, 0.83, 4.54, 2.12, 0.25, "地址/同步随请求变化", 10.5, True, RED, PP_ALIGN.CENTER)
    add_textbox(slide, 1.22, 5.02, 1.48, 0.30, "Capture / Replay 不稳定", 11, True, RED, PP_ALIGN.CENTER)

    # Center transformation.
    add_textbox(slide, 3.70, 4.12, 1.18, 0.32, "状态\n虚拟化", 13.5, True, ORANGE, PP_ALIGN.CENTER)
    add_arrow(slide, 3.52, 4.25, 4.70, 4.25, ORANGE, 3.0)

    # SEW architecture.
    add_rect(slide, 4.88, 3.20, 5.56, 2.25, LIGHT_GREEN, RGBColor(0x8A, 0xC7, 0xAE), radius=True)
    add_textbox(slide, 5.04, 3.28, 5.23, 0.28, "SEW-Offload：重放安全的专家状态管理运行时", 12.5, True, GREEN, PP_ALIGN.CENTER)

    add_pill(slide, 5.08, 3.76, 0.86, 0.34, "Router", PALE_BLUE, RGBColor(0xB9, 0xC7, 0xD8), 9.5)
    add_arrow(slide, 5.93, 3.93, 6.34, 3.93, GREEN, 2.2)
    add_pill(slide, 6.34, 3.72, 1.45, 0.43, "Staging\nController", LIGHT_ORANGE, RGBColor(0xF0, 0xC5, 0x7A), 9.2)
    add_arrow(slide, 7.78, 3.93, 8.18, 3.93, GREEN, 2.2)

    add_rect(slide, 8.18, 3.58, 1.92, 0.78, WHITE, RGBColor(0xA9, 0xC4, 0xB8), radius=True)
    add_textbox(slide, 8.25, 3.63, 1.78, 0.22, "Fixed NPU Slots", 9.5, True, GREEN, PP_ALIGN.CENTER)
    for i in range(4):
        add_rect(slide, 8.34 + i * 0.40, 3.90, 0.30, 0.26, RGBColor(0x78, 0xB7, 0x9D), RGBColor(0x78, 0xB7, 0x9D), radius=True)

    add_rect(slide, 5.42, 4.52, 1.75, 0.46, WHITE, RGBColor(0xA9, 0xC4, 0xB8), radius=True)
    add_textbox(slide, 5.50, 4.58, 1.58, 0.20, "CPU Expert Store", 9.5, True, GREEN, PP_ALIGN.CENTER)
    add_arrow(slide, 6.32, 4.50, 6.92, 4.15, GREEN, 1.8)

    add_rect(slide, 7.42, 4.56, 1.26, 0.38, WHITE, RGBColor(0xA9, 0xC4, 0xB8), radius=True)
    add_textbox(slide, 7.46, 4.61, 1.18, 0.18, "log2phy", 9.5, True, GREEN, PP_ALIGN.CENTER)
    add_arrow(slide, 7.85, 4.56, 8.30, 4.24, GREEN, 1.8)

    add_rect(slide, 8.93, 4.58, 1.28, 0.42, RGBColor(0xDB, 0xF0, 0xE7), GREEN, radius=True)
    add_textbox(slide, 8.96, 4.62, 1.22, 0.22, "ACLGraph\nReplay", 8.8, True, GREEN, PP_ALIGN.CENTER)
    add_arrow(slide, 9.60, 4.36, 9.60, 4.58, GREEN, 1.8)

    add_pill(slide, 5.08, 5.02, 1.62, 0.25, "1 路由先于重放完成", WHITE, RGBColor(0xA9, 0xC4, 0xB8), 8.5, True, GREEN)
    add_pill(slide, 6.90, 5.02, 1.54, 0.25, "2 Slot 生命周期保护", WHITE, RGBColor(0xA9, 0xC4, 0xB8), 8.5, True, GREEN)
    add_pill(slide, 8.62, 5.02, 1.50, 0.25, "3 B2 分波不丢专家", WHITE, RGBColor(0xA9, 0xC4, 0xB8), 8.5, True, GREEN)

    # Result cards.
    add_card(slide, 10.68, 3.20, 2.20, 0.98, "图执行", "Capture ✓", "ACLGraph 完成重放", PALE_BLUE, BLUE)
    add_card(slide, 10.68, 4.30, 2.20, 0.98, "显存状态", "32 Slots", "Slot Bank 3.4 GiB", LIGHT_GREEN, GREEN)
    add_card(slide, 10.68, 5.40, 2.20, 0.98, "Prefill 溢出", "126 → 4波", "active experts > slots", LIGHT_ORANGE, ORANGE)

    add_band(
        slide,
        6.50,
        "当前效果",
        "ShareGPT mixed / Qwen3-30B-A3B 初步运行：200/200 请求成功，P50 TTFT 1.37s，P50 TPOT 83.6ms，输出 10.5 tok/s",
        RED,
        PALE_RED,
        h=0.55,
        body_size=16,
    )
    add_textbox(
        slide,
        8.78,
        7.08,
        4.10,
        0.18,
        "注：数值来自 2026-07-06 单次初步运行；正式汇报可替换为多次重复与基线对比结果",
        6.5,
        False,
        MID,
        PP_ALIGN.RIGHT,
    )

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
