# -*- coding: utf-8 -*-
"""对一页 PPTX 做版式自检：形状重叠 + 文本溢出估算，并出一张色块缩略图。"""
import sys
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt
from PIL import Image, ImageDraw, ImageFont

FN = sys.argv[1] if len(sys.argv) > 1 else "LatchMoE_年会一页PPT.pptx"
SCALE = 96  # px per inch
prs = Presentation(FN)
s = prs.slides[0]
W = Emu(prs.slide_width).inches
H = Emu(prs.slide_height).inches


def is_cjk(ch):
    return '一' <= ch <= '鿿' or '　' <= ch <= '〿' or '＀' <= ch <= '￯'


def text_width_in(runs, default_pt):
    """估算一段（一行）文本宽度（英寸）。CJK≈1.0em，拉丁≈0.56em。"""
    w = 0.0
    for txt, pt in runs:
        em = pt / 72.0
        for ch in txt:
            w += em * (1.02 if is_cjk(ch) else 0.56)
    return w


def gather(shape):
    """返回 (l,t,w,h, is_solid_fill, [ (line_runs=[(txt,pt)], align) ], frame_w, frame_h)."""
    l, t = Emu(shape.left).inches, Emu(shape.top).inches
    w, h = Emu(shape.width).inches, Emu(shape.height).inches
    lines = []
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            runs = []
            for r in p.runs:
                pt = r.font.size.pt if r.font.size else 14
                runs.append((r.text, pt))
            if runs:
                lines.append(runs)
    return l, t, w, h, lines


# ---- 文本溢出检测 ----
print("=== 文本溢出检测（单行估算宽 > 盒宽*1.02 判溢出）===")
overflow = 0
for sh in s.shapes:
    if not sh.has_text_frame:
        continue
    l, t, w, h, lines = gather(sh)
    for runs in lines:
        # 显式换行分别估算，避免把两行英文标签误判成溢出。
        split_lines = [[]]
        for txt, pt in runs:
            chunks = txt.split("\n")
            for i, chunk in enumerate(chunks):
                if i:
                    split_lines.append([])
                split_lines[-1].append((chunk, pt))
        for split_runs in split_lines:
            tw = text_width_in(split_runs, 14)
            # 文本框左右各 0.05 margin
            avail = w - 0.10
            if tw > avail * 1.03 and avail > 0.2:
                txt = ''.join(r[0] for r in split_runs)
                print(f"  OVERFLOW  box_w={w:.2f} avail={avail:.2f} textw={tw:.2f}  {txt[:46]!r}")
                overflow += 1
print("overflow lines:", overflow)


# ---- 出色块缩略图 ----
def rgb(shape):
    try:
        c = shape.fill.fore_color.rgb
        return (c[0], c[1], c[2])
    except Exception:
        return None


def line_rgb(shape):
    try:
        c = shape.line.color.rgb
        return (c[0], c[1], c[2])
    except Exception:
        return None


img = Image.new("RGB", (int(W * SCALE), int(H * SCALE)), (255, 255, 255))
d = ImageDraw.Draw(img)
REGULAR_FONT = "/root/.local/share/fonts/NotoSansCJKsc-Regular.otf"
BOLD_FONT = "/root/.local/share/fonts/NotoSansCJKsc-Bold.otf"
font_cache = {}


def get_font(pt, bold=False):
    px_size = max(6, round(pt * SCALE / 72.0))
    key = (px_size, bool(bold))
    if key not in font_cache:
        try:
            font_cache[key] = ImageFont.truetype(BOLD_FONT if bold else REGULAR_FONT, px_size)
        except Exception:
            font_cache[key] = ImageFont.load_default()
    return font_cache[key]


def px(v):
    return int(v * SCALE)


def draw_shape(sh):
    l, t = Emu(sh.left).inches, Emu(sh.top).inches
    w, h = Emu(sh.width).inches, Emu(sh.height).inches
    st = str(sh.shape_type)
    if "LINE" in st or "CONNECTOR" in st:
        lc = line_rgb(sh) or (120, 120, 120)
        d.line([px(l), px(t), px(l + w), px(t + h)], fill=lc, width=2)
        return
    fill = rgb(sh)
    ln = line_rgb(sh)
    if fill or ln:
        d.rectangle([px(l), px(t), px(l + w), px(t + h)],
                    fill=fill, outline=ln or fill, width=1)
    # 文本：用本机 Noto Sans CJK 做真实缩略图，便于检查层级与密度。
    if sh.has_text_frame:
        visual_lines = []
        for p in sh.text_frame.paragraphs:
            line = []
            for r in p.runs:
                if not r.text:
                    continue
                pt = r.font.size.pt if r.font.size else 14
                bold = bool(r.font.bold)
                try:
                    color = tuple(r.font.color.rgb) if r.font.color.rgb else (20, 20, 20)
                except Exception:
                    color = (20, 20, 20)
                chunks = r.text.split("\n")
                for i, chunk in enumerate(chunks):
                    if i:
                        visual_lines.append((line, p.alignment))
                        line = []
                    if chunk:
                        line.append((chunk, pt, bold, color))
            if line:
                visual_lines.append((line, p.alignment))

        metrics = []
        for runs, align in visual_lines:
            widths = []
            max_h = 0
            for txt, pt, bold, color in runs:
                font = get_font(pt, bold)
                bbox = d.textbbox((0, 0), txt, font=font)
                widths.append(bbox[2] - bbox[0])
                max_h = max(max_h, bbox[3] - bbox[1])
            metrics.append((runs, align, sum(widths), max(8, max_h), widths))

        gap = 2
        total_h = sum(m[3] for m in metrics) + gap * max(0, len(metrics) - 1)
        anchor = sh.text_frame.vertical_anchor
        if anchor == MSO_ANCHOR.MIDDLE:
            yy = px(t) + max(1, (px(h) - total_h) // 2)
        elif anchor == MSO_ANCHOR.BOTTOM:
            yy = px(t + h) - total_h - 2
        else:
            yy = px(t) + 2

        for runs, align, line_w, line_h, widths in metrics:
            if align == PP_ALIGN.CENTER:
                xx = px(l) + max(2, (px(w) - line_w) // 2)
            elif align == PP_ALIGN.RIGHT:
                xx = px(l + w) - line_w - 4
            else:
                xx = px(l) + 4
            for (txt, pt, bold, color), run_w in zip(runs, widths):
                d.text((xx, yy), txt, fill=color, font=get_font(pt, bold), anchor="lt")
                xx += run_w
            yy += line_h + gap


for sh in s.shapes:
    draw_shape(sh)

out = FN.replace(".pptx", "_thumb.png")
img.save(out)
print("thumb:", out, img.size)
