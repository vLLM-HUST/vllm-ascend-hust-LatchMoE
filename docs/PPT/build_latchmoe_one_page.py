# -*- coding: utf-8 -*-
"""年会一页 PPT：LatchMoE —— 昇腾 MoE 专家卸载 × ACLGraph 图重放协同。

版式沿用课题组既有材料的蓝色标题、红色问题带和底部结果条；主体用可编辑
矢量形状重画“动态路由 vs. 稳定重放”的冲突，并用三项机制闭环回答挑战。
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("LatchMoE_年会一页PPT.pptx")
FONT = "Microsoft YaHei"

# ---- 组内模板配色（取自 系数量化协同.png / KV前缀复用.png / BidKV.pptx）----
NAVY       = RGBColor(0x11, 0x3A, 0x88)   # 标题
DARKRED    = RGBColor(0xC0, 0x00, 0x00)   # 深红标题条 / tab
PALE_RED   = RGBColor(0xF9, 0xE5, 0xE5)   # 粉红正文底
PALE_BLUE  = RGBColor(0xE5, 0xEA, 0xF5)   # 浅蓝底
BLUE       = RGBColor(0x03, 0x39, 0xC4)   # 蓝色强调
GREEN      = RGBColor(0x2E, 0x8B, 0x6B)   # 绿点缀
LT_GREEN   = RGBColor(0xF1, 0xFE, 0xF5)
TEAL_LINE  = RGBColor(0xA9, 0xC4, 0xB8)
ORANGE     = RGBColor(0xE0, 0x8A, 0x1A)
LT_ORANGE  = RGBColor(0xFD, 0xF2, 0xDF)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
MID        = RGBColor(0x5C, 0x6B, 0x7A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LN    = RGBColor(0xD7, 0xDF, 0xEA)
RED_LN     = RGBColor(0xE0, 0xA8, 0xA8)


def _noshadow(shp):
    try:
        shp.shadow.inherit = False
    except Exception:
        pass


def rect(slide, x, y, w, h, fill, line=None, radius=0.0, lw=1.0):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    if radius:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    _noshadow(s)
    return s


def text(slide, x, y, w, h, runs, size=14, bold=False, color=DARK,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.0, wrap=True):
    """runs 可为字符串，或 [(片段, {size,bold,color}), ...] 富文本列表。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    if isinstance(runs, str):
        runs = [runs]
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        frags = line if isinstance(line, list) else [line]
        for frag in frags:
            if isinstance(frag, tuple):
                txt, st = frag
            else:
                txt, st = frag, {}
            r = p.add_run(); r.text = txt
            r.font.name = st.get("font", FONT)
            r.font.size = Pt(st.get("size", size))
            r.font.bold = st.get("bold", bold)
            r.font.color.rgb = st.get("color", color)
    return box


def tri(slide, x, y, w, h, fill, flip=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if flip:
        s.rotation = 0
    _noshadow(s)
    return s


def arrow(slide, x1, y1, x2, y2, color=MID, width=2.0, dash=False, both=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    c.line.end_arrowhead = True
    if both:
        c.line.begin_arrowhead = True
    if dash:
        c.line.dash_style = 4
    _noshadow(c)
    return c


def pill(slide, x, y, w, h, txt, fill, line=None, size=11, bold=True, color=DARK):
    rect(slide, x, y, w, h, fill, line, radius=0.5, lw=1.0)
    text(slide, x, y, w, h, txt, size, bold, color, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def tab(slide, y, label, body, x=0.20, w=6.30, h=0.32, label_w=1.55, body_size=11.5):
    """深红 tab + 细体标题行（对齐 BidKV 的『背景：…』『挑战：…』）。"""
    rect(slide, x, y, label_w, h, DARKRED, radius=0.0)
    text(slide, x, y, label_w, h, label, 12.5, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    text(slide, x + label_w + 0.12, y, w - label_w - 0.12, h, body, body_size, True, DARKRED,
         PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def metric_card(slide, x, y, w, h, name, bval, bunit, aval, aunit, mult, accent):
    """紧凑横排成绩卡：指标名条 + before →(倍率) after。"""
    rect(slide, x, y, w, h, WHITE, GREY_LN, radius=0.10, lw=1.0)
    rect(slide, x + 0.02, y + 0.02, 0.10, h - 0.04, accent, radius=0.0)
    text(slide, x + 0.20, y + 0.09, w - 0.30, 0.26, name, 12, True, NAVY, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    row_y = y + 0.44
    # before
    text(slide, x + 0.16, row_y, 1.16, 0.30, bval, 17, True, MID, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    text(slide, x + 0.16, row_y + 0.30, 1.16, 0.18, "单算子 " + bunit, 8.5, False, MID, PP_ALIGN.CENTER)
    # arrow
    arrow(slide, x + 1.40, row_y + 0.22, x + 1.86, row_y + 0.22, accent, 2.6)
    pill(slide, x + 1.20, row_y - 0.20, 0.86, 0.24, mult, accent, accent, 11, True, WHITE)
    # after
    text(slide, x + 1.94, row_y - 0.06, w - 2.06, 0.44, aval, 26, True, accent, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    text(slide, x + 1.94, row_y + 0.40, w - 2.06, 0.18, "LatchMoE " + aunit, 8.5, True, accent, PP_ALIGN.CENTER)


def expert_sq(slide, x, y, s, fill, line=None):
    rect(slide, x, y, s, s, fill, line or fill, radius=0.18, lw=0.75)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 13.333, 7.5, WHITE)

    # ===== 标题 =====
    text(slide, 0.24, 0.16, 12.85, 0.52,
         "课题：LatchMoE——昇腾 MoE 专家卸载与 ACLGraph 协同",
         27, True, NAVY, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    text(slide, 0.26, 0.70, 12.8, 0.26,
         "把“动态换专家”锁在图外，把“固定地址计算”留在图内",
         12.5, False, MID, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    rect(slide, 0.26, 0.66, 3.35, 0.028, DARKRED)

    # ===== 核心挑战带（全宽）=====
    by = 1.06
    rect(slide, 0.20, by, 12.92, 0.60, PALE_RED, RED_LN, radius=0.06)
    rect(slide, 0.20, by, 1.72, 0.60, DARKRED, radius=0.06)
    tri(slide, 1.72, by, 0.26, 0.60, DARKRED)
    text(slide, 0.20, by, 1.72, 0.60, "核心挑战", 14, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    text(slide, 2.14, by, 10.85, 0.60,
         [[("路由每步变化，ACLGraph 却要求地址与依赖稳定", {"bold": True, "color": DARKRED, "size": 14}),
           ("：传统按需搬移只能退回单算子模式，吞吐仅 ", {"color": DARK, "size": 12.5}),
           ("4.71 tok/s", {"bold": True, "color": DARKRED, "size": 13.5}),
           ("。", {"color": DARK, "size": 12.5})]],
         12.5, False, DARK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    # ===== 左列：背景 =====
    lx, lw = 0.20, 6.28
    tab(slide, 1.86, "背景", "专家权重超出 HBM，必须在 Host 与 NPU 间按需搬移", x=lx, w=lw, h=0.32)
    # 背景示意图
    bg_y = 2.28
    rect(slide, lx + 0.06, bg_y, 2.62, 1.16, PALE_BLUE, GREY_LN, radius=0.06)
    text(slide, lx + 0.10, bg_y + 0.04, 2.54, 0.20, "HBM 容量受限 → 冷专家驻留 Host", 9, True, NAVY, PP_ALIGN.CENTER)
    rect(slide, lx + 0.18, bg_y + 0.30, 1.02, 0.74, WHITE, BLUE, radius=0.06)
    text(slide, lx + 0.18, bg_y + 0.31, 1.02, 0.18, "NPU HBM", 8, True, BLUE, PP_ALIGN.CENTER)
    for i in range(2):
        for j in range(3):
            expert_sq(slide, lx + 0.28 + j * 0.28, bg_y + 0.54 + i * 0.24, 0.20, PALE_BLUE, BLUE)
    rect(slide, lx + 1.52, bg_y + 0.28, 1.02, 0.78, LT_ORANGE, ORANGE, radius=0.06)
    text(slide, lx + 1.52, bg_y + 0.30, 1.02, 0.16, "Host 内存", 8, True, ORANGE, PP_ALIGN.CENTER)
    for i in range(3):
        for j in range(4):
            expert_sq(slide, lx + 1.58 + j * 0.22, bg_y + 0.50 + i * 0.18, 0.16, WHITE, ORANGE)
    arrow(slide, lx + 1.52, bg_y + 0.66, lx + 1.24, bg_y + 0.66, ORANGE, 2.0, both=True)
    text(slide, lx + 2.85, bg_y - 0.02, 3.55, 1.22,
         [[("Qwen3-30B-A3B 的专家权重无法全部常驻有限 HBM，", {"size": 10.5})],
          [("需从 Host 按路由结果调入。", {"size": 10.5})],
          [("但 ACLGraph 重放要求算子的", {"size": 10.5})],
          [("地址 / 形状 / 依赖保持稳定", {"bold": True, "color": NAVY, "size": 10.5}), ("。", {"size": 10.5})]],
         11, False, DARK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, line_spacing=1.12)

    # ===== 左列：挑战 =====
    tab(slide, 3.60, "冲突", "动态搬移进入捕获区，图重放失效；放在图外又要保证数据就绪", x=lx, w=lw, h=0.32)
    ch_y = 4.02
    rect(slide, lx + 0.06, ch_y, 2.62, 1.28, RGBColor(0xFF, 0xF6, 0xF6), RED_LN, radius=0.06)
    text(slide, lx + 0.10, ch_y + 0.04, 2.54, 0.18, "逐步变化的路由撞上固定图边界", 9, True, DARKRED, PP_ALIGN.CENTER)
    pill(slide, lx + 0.20, ch_y + 0.34, 0.80, 0.30, "Router", PALE_BLUE, GREY_LN, 9, True, NAVY)
    arrow(slide, lx + 1.00, ch_y + 0.49, lx + 1.28, ch_y + 0.49, DARKRED, 2.0)
    pill(slide, lx + 1.28, ch_y + 0.34, 1.02, 0.30, "H2D 搬移", PALE_RED, RED_LN, 9, True, DARKRED)
    # graph boundary
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(lx + 1.14), Inches(ch_y + 0.28), Inches(lx + 1.14), Inches(ch_y + 1.18))
    line.line.color.rgb = DARKRED; line.line.width = Pt(1.6); line.line.dash_style = 4
    _noshadow(line)
    rect(slide, lx + 0.40, ch_y + 0.74, 1.90, 0.28, WHITE, RED_LN, radius=0.10)
    text(slide, lx + 0.40, ch_y + 0.74, 1.90, 0.28, "地址 / 同步随请求变", 9, True, DARKRED, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    text(slide, lx + 0.10, ch_y + 1.04, 2.54, 0.22, "Capture / Replay 失效", 9.5, True, DARKRED, PP_ALIGN.CENTER)
    text(slide, lx + 2.85, ch_y - 0.02, 3.55, 1.32,
         [[("① 动态专家如何绑定到", {"size": 10.5}),
           ("稳定接口", {"bold": True, "color": DARKRED, "size": 11}), ("？", {"size": 10.5})],
          [("② Prefill 活跃专家超过 Slot 容量怎么办？", {"size": 10.5})],
          [("③ 传输流写 Slot 与计算流读 Slot 如何避免竞争？", {"size": 10.5})],
          [("目标：", {"bold": True, "color": DARKRED, "size": 10.5}),
           ("内容可变、地址不变，且读写安全。", {"bold": True, "color": DARKRED, "size": 10.5})]],
         10.5, False, DARK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, line_spacing=1.1)

    # ===== 右列：方法 =====
    rx, rw = 6.72, 6.40
    rect(slide, rx, 1.86, rw, 3.48, LT_GREEN, TEAL_LINE, radius=0.03)
    rect(slide, rx, 1.86, rw, 0.38, GREEN, radius=0.0)
    text(slide, rx, 1.86, rw, 0.38, "方法：LatchMoE——用固定 Slot 解耦动态路由与图重放",
         13, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    # ---- 流程图 ----
    fy = 2.42
    pill(slide, rx + 0.22, fy + 0.30, 0.84, 0.34, "Router", PALE_BLUE, GREY_LN, 10, True, NAVY)
    arrow(slide, rx + 1.06, fy + 0.47, rx + 1.40, fy + 0.47, GREEN, 2.2)
    rect(slide, rx + 1.40, fy + 0.22, 1.30, 0.50, LT_ORANGE, ORANGE, radius=0.10)
    text(slide, rx + 1.40, fy + 0.22, 1.30, 0.50, "Staging\nController", 9.5, True, ORANGE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    arrow(slide, rx + 2.70, fy + 0.47, rx + 3.06, fy + 0.47, GREEN, 2.2)
    # fixed slots
    rect(slide, rx + 3.06, fy + 0.10, 1.66, 0.74, WHITE, GREEN, radius=0.08)
    text(slide, rx + 3.06, fy + 0.12, 1.66, 0.18, "Fixed NPU Slots", 9, True, GREEN, PP_ALIGN.CENTER)
    for i in range(4):
        expert_sq(slide, rx + 3.20 + i * 0.36, fy + 0.42, 0.28, RGBColor(0x78, 0xB7, 0x9D), GREEN)
    arrow(slide, rx + 4.72, fy + 0.47, rx + 5.08, fy + 0.47, GREEN, 2.2)
    rect(slide, rx + 5.08, fy + 0.22, 1.12, 0.50, RGBColor(0xDB, 0xF0, 0xE7), GREEN, radius=0.10)
    text(slide, rx + 5.08, fy + 0.22, 1.12, 0.50, "ACLGraph\nReplay ✓", 9.5, True, GREEN, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    # host store feeding staging
    rect(slide, rx + 1.28, fy + 0.94, 1.54, 0.34, WHITE, TEAL_LINE, radius=0.10)
    text(slide, rx + 1.28, fy + 0.94, 1.54, 0.34, "Host Expert Store", 9, True, GREEN, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    arrow(slide, rx + 2.05, fy + 0.94, rx + 2.05, fy + 0.72, GREEN, 1.8)
    # log2phy label under slots
    rect(slide, rx + 3.40, fy + 0.94, 1.10, 0.32, WHITE, TEAL_LINE, radius=0.12)
    text(slide, rx + 3.40, fy + 0.94, 1.10, 0.32, "log2phy 映射", 8.5, True, GREEN, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    arrow(slide, rx + 3.95, fy + 0.94, rx + 3.95, fy + 0.84, GREEN, 1.8)

    # ---- 三个核心机制（两行式：加粗标题 + 灰色说明）----
    mech = [
        ("1", "固定 Slot + 持久 log2phy",
         "图外换内容与更新映射；图内 Grouped MLP 只读固定地址"),
        ("2", "容量受限 B2 波次 Prefill",
         "活跃专家超过槽位时分波调入与计算，保持完整模型语义"),
        ("3", "COMPUTING 保护的 Slot 生命周期",
         "事件同步保护在算 Slot，避免传输流写与计算流读竞争"),
    ]
    myy = fy + 1.40
    for i, (n, t, d) in enumerate(mech):
        yy = myy + i * 0.52
        rect(slide, rx + 0.22, yy + 0.02, 0.32, 0.32, GREEN, radius=0.5)
        text(slide, rx + 0.22, yy + 0.02, 0.32, 0.32, n, 12.5, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        text(slide, rx + 0.66, yy - 0.02, 5.60, 0.22, t, 11.5, True, GREEN, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        text(slide, rx + 0.66, yy + 0.19, 5.62, 0.20, d, 10, False, MID, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    # ===== 达成效果带 =====
    ey = 5.56
    rect(slide, 0.20, ey, 12.92, 0.40, PALE_RED, RED_LN, radius=0.06)
    rect(slide, 0.20, ey, 1.72, 0.40, DARKRED, radius=0.06)
    tri(slide, 1.72, ey, 0.22, 0.40, DARKRED)
    text(slide, 0.20, ey, 1.72, 0.40, "达成效果", 13.5, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    text(slide, 2.12, ey, 10.9, 0.40,
         [[("ShareGPT · Qwen3-30B-A3B：相对优化前单算子模式，LatchMoE 吞吐提升 ", {"size": 12}),
           ("2.70×", {"bold": True, "color": DARKRED, "size": 14}),
           ("，TTFT / TPOT 分别下降 58% / 63%。", {"size": 12})]],
         12, False, DARK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    # ===== 三张成绩卡 =====
    cy, ch, gap = 6.06, 1.20, 0.22
    cw = (12.92 - 2 * gap) / 3
    metric_card(slide, 0.20 + 0 * (cw + gap), cy, cw, ch, "首 Token 延迟 TTFT",
                "2283.2", "ms", "952.3", "ms", "↓58%", BLUE)
    metric_card(slide, 0.20 + 1 * (cw + gap), cy, cw, ch, "单步延迟 TPOT",
                "192.5", "ms", "70.57", "ms", "↓63%", ORANGE)
    metric_card(slide, 0.20 + 2 * (cw + gap), cy, cw, ch, "端到端吞吐",
                "4.71", "tok/s", "12.7", "tok/s", "2.70×", GREEN)

    text(slide, 0.20, 7.30, 12.92, 0.18,
         "对比口径：同一 ShareGPT 数据集、同一 Qwen3-30B-A3B 模型；优化前为单算子模式，优化后为 LatchMoE。",
         8, False, MID, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    prs.save(OUT)
    print("saved:", OUT)


if __name__ == "__main__":
    build()
